"""
Build DEMO_WALKTHROUGH.pptx — the four live-demo moments from the playbook (§9),
each explained in mild detail with code snippets + flow diagrams.

Regenerate with:   py -3 docs/build_demo_deck.py
Requires:          py -3 -m pip install python-pptx

Theme: TD Bank colours (green + gold on white). 16:9.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (TD Bank colours) -----------------------------------------------
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x00, 0x8A, 0x00)
GREEN_DARK = RGBColor(0x00, 0x6B, 0x00)
GREEN_DEEP = RGBColor(0x00, 0x4F, 0x00)
YELLOW     = RGBColor(0xFF, 0xC7, 0x2C)
BG      = WHITE
PANEL   = RGBColor(0xF2, 0xF4, 0xF5)
PANEL2  = RGBColor(0xE8, 0xF5, 0xE8)
ACCENT  = GREEN
ACCENT2 = GREEN_DARK
AMBER   = RGBColor(0xB4, 0x53, 0x09)
RED     = RGBColor(0xC7, 0x10, 0x2E)
TEXT    = RGBColor(0x1C, 0x1C, 0x1C)
MUTED   = RGBColor(0x55, 0x5E, 0x66)
CODE    = RGBColor(0xE8, 0xF5, 0xE8)
CODE_BG = RGBColor(0x06, 0x2A, 0x14)

TITLE_FONT = "Segoe UI Semibold"
BODY_FONT  = "Segoe UI"
MONO_FONT  = "Consolas"
SW, SH = Emu(12192000), Emu(6858000)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s


def _set(run, size, color, font=BODY_FONT, bold=False, italic=False):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.name = font; run.font.bold = bold; run.font.italic = italic


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, Pt(2))
    return tf


def para(tf, text, size, color, font=BODY_FONT, bold=False, italic=False,
         space_after=6, level=0, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.level = level; p.alignment = align; p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    _set(r, size, color, font, bold, italic)
    return p


def header(s, kicker, title):
    tf = textbox(s, Inches(0.6), Inches(0.30), Inches(12.1), Inches(0.3))
    para(tf, kicker.upper(), 12, ACCENT, BODY_FONT, bold=True, first=True, space_after=0)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.62), Inches(0.9), Pt(5))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
    tf2 = textbox(s, Inches(0.6), Inches(0.74), Inches(12.1), Inches(0.95))
    para(tf2, title, 28, GREEN_DEEP, TITLE_FONT, bold=True, first=True, space_after=0)


def footer(s, n):
    tf = textbox(s, Inches(0.6), Inches(7.05), Inches(9.0), Inches(0.35))
    r = tf.paragraphs[0].add_run(); r.text = "Banking Admin Portal  ·  Live demo walkthrough"; _set(r, 9, MUTED)
    tfn = textbox(s, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.35))
    pn = tfn.paragraphs[0]; pn.alignment = PP_ALIGN.RIGHT
    rn = pn.add_run(); rn.text = str(n); _set(rn, 9, MUTED)


def panel(s, x, y, w, h, fill=PANEL, line=None, radius=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(1.25)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def box(s, x, y, w, h, title, subs=(), fill=PANEL2, edge=ACCENT2,
        title_color=TEXT, body_color=MUTED, fs_title=12, fs_body=9.5):
    shp = panel(s, x, y, w, h, fill=fill, line=edge)
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(5); tf.margin_right = Pt(5)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; _set(r, fs_title, title_color, BODY_FONT, bold=True)
    for sub in subs:
        sp = tf.add_paragraph(); sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run(); sr.text = sub; _set(sr, fs_body, body_color, MONO_FONT)
    return shp


def arrow(s, x1, y1, x2, y2, color=ACCENT, width=2.0):
    c = s.shapes.add_connector(2, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = Pt(width); c.shadow.inherit = False
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def code_panel(s, x, y, w, h, code, fs=12, caption=None):
    panel(s, x, y, w, h, fill=CODE_BG, line=GREEN_DARK)
    tf = textbox(s, x + Pt(10), y + Pt(8), w - Pt(20), h - Pt(16))
    if caption:
        p = tf.paragraphs[0]; r = p.add_run(); r.text = caption
        _set(r, fs - 1, YELLOW, MONO_FONT, bold=True); p.space_after = Pt(6)
    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if (i == 0 and not caption) else tf.add_paragraph()
        p.space_after = Pt(0); p.line_spacing = 1.05
        r = p.add_run(); r.text = line if line else " "
        _set(r, fs, CODE, MONO_FONT)


def bullets(s, x, y, w, h, items, fs=15, gap=8):
    tf = textbox(s, x, y, w, h)
    for i, it in enumerate(items):
        text, level, color, bold = (it + (0, TEXT, False))[:4]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level; p.space_after = Pt(gap)
        r = p.add_run(); r.text = ("▸ " if level == 0 else "– ") + text
        _set(r, fs - level, color, BODY_FONT, bold=bold)
    return tf


def whatwhy(s, what, why, y=Inches(1.75)):
    """Two side-by-side panels: What to show / Why it lands."""
    panel(s, Inches(0.6), y, Inches(6.0), Inches(1.5), fill=PANEL2, line=ACCENT)
    tf = textbox(s, Inches(0.85), y + Pt(10), Inches(5.5), Inches(1.2))
    para(tf, "WHAT TO SHOW", 12, GREEN_DEEP, BODY_FONT, bold=True, first=True, space_after=5)
    para(tf, what, 13, TEXT)
    panel(s, Inches(6.85), y, Inches(5.85), Inches(1.5), fill=PANEL, line=YELLOW)
    tf2 = textbox(s, Inches(7.1), y + Pt(10), Inches(5.4), Inches(1.2))
    para(tf2, "WHY IT LANDS", 12, AMBER, BODY_FONT, bold=True, first=True, space_after=5)
    para(tf2, why, 13, TEXT)


# =============================================================================
# 1 — TITLE
# =============================================================================
s = slide(bg=GREEN)
rib = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.55), SW, Pt(5))
rib.fill.solid(); rib.fill.fore_color.rgb = YELLOW; rib.line.fill.background(); rib.shadow.inherit = False
tf = textbox(s, Inches(0.9), Inches(1.45), Inches(11.5), Inches(1.0))
para(tf, "BANKING ADMIN PORTAL", 16, YELLOW, BODY_FONT, bold=True, first=True, space_after=2)
tf2 = textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.4))
para(tf2, "Live demo — four moments that land", 38, WHITE, TITLE_FONT, bold=True, first=True, space_after=4)
para(tf2, "Audit log in real time · accessibility you can't fake · the test net · end-to-end tracing", 17, PANEL2)
tf3 = textbox(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.8))
para(tf3, "Each moment: what to click, why it lands, and the code behind it", 14, YELLOW, BODY_FONT, bold=True, first=True)

# =============================================================================
# 2 — OVERVIEW
# =============================================================================
s = slide(); header(s, "The plan", "Four moments, in order")
cards = [
    ("1 · Audit log, live", "Toggle a status → the UPDATE diff appears at the top of the audit log in real time.", ACCENT),
    ("2 · A11y you can't fake", "Delete → Escape → focus returns to the exact button. Keyboard + screen-reader correct.", AMBER),
    ("3 · The test net", "Open the Cypress runner, pick a spec, let it walk the UI. Don't talk over it.", ACCENT),
    ("4 · End-to-end tracing", "One X-Correlation-Id flows browser → API → log. Zero extra dependencies.", AMBER),
]
y = Inches(1.85)
for i, (t, d, c) in enumerate(cards):
    yy = y + Inches(1.28) * (i // 2)
    xx = Inches(0.6) + Inches(6.25) * (i % 2)
    panel(s, xx, yy, Inches(6.0), Inches(1.12), fill=PANEL, line=c)
    tf = textbox(s, xx + Pt(12), yy + Pt(8), Inches(5.6), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, 15, GREEN_DEEP, BODY_FONT, bold=True, first=True, space_after=3)
    para(tf, d, 11.5, MUTED)
para(textbox(s, Inches(0.6), Inches(6.4), Inches(12.0), Inches(0.5)),
     "Golden rule: show the behaviour first, then open the code only if they ask how.",
     13, AMBER, italic=True, first=True)
footer(s, 2)

# =============================================================================
# 3 — MOMENT 1 : what/why + flow
# =============================================================================
s = slide(); header(s, "Moment 1 · Audit log", "Toggle status → the diff appears, live")
whatwhy(s,
        "On an employee, click \"Mark INACTIVE\". Scroll to the audit log — a new UPDATE entry is at the top.",
        "Shows the append-only audit working in real time, with a clean field-level diff. The banking compliance story, visible.")
# flow diagram
y = Inches(3.7)
nodes = [("Toggle status", ["component"]), ("PATCH /employees/:id", ["{ status }"]),
         ("service.patch()", ["before → update", "→ audit.diff()"]), ("Audit entry", ["at top of log"])]
xs = [Inches(0.6), Inches(3.65), Inches(6.85), Inches(10.25)]
ws = [Inches(2.7), Inches(2.9), Inches(3.1), Inches(2.45)]
for (t, subs), x, w in zip(nodes, xs, ws):
    box(s, x, y, w, Inches(1.15), t, subs, edge=ACCENT, fs_title=12.5)
for i in range(3):
    arrow(s, xs[i] + ws[i], y + Inches(0.57), xs[i + 1], y + Inches(0.57))
code_panel(s, Inches(0.6), Inches(5.25), Inches(12.1), Inches(1.45),
           "{ \"action\": \"UPDATE\", \"entity\": \"Employee\",\n"
           "  \"changes\": [ { \"field\": \"status\", \"before\": \"ACTIVE\", \"after\": \"INACTIVE\" } ],\n"
           "  \"actor\": \"admin\", \"correlationId\": \"…\" }", fs=12, caption="// the audit entry the toggle produces")
footer(s, 3)

# =============================================================================
# 4 — MOMENT 1 : code
# =============================================================================
s = slide(); header(s, "Moment 1 · Code", "From click to audit entry")
code_panel(s, Inches(0.6), Inches(1.75), Inches(6.0), Inches(2.5),
           "// employee-detail.component.ts\n"
           "toggleStatus(emp: Employee): void {\n"
           "  this.facade.patchStatus(\n"
           "    emp.employeeId,\n"
           "    emp.status === 'ACTIVE'\n"
           "      ? 'INACTIVE' : 'ACTIVE');\n"
           "}", fs=12.5)
code_panel(s, Inches(6.85), Inches(1.75), Inches(5.85), Inches(2.5),
           "// employee.service.js — patch()\n"
           "const before = repo.findById(id);   // snapshot\n"
           "const after  = repo.update(id, patch);\n"
           "AuditService.recordEmployeeUpdated(\n"
           "  before, after, context);          // diff\n"
           "return after;", fs=12.5)
bullets(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.0), [
    ("The component only dispatches through the facade — it never mutates state directly.", 0, TEXT, False),
    ("The service snapshots the row BEFORE the write, so the audit diff is exact.", 0, TEXT, False),
    ("audit.diff() emits only changed fields — a no-op write records nothing.", 0, ACCENT, True),
], fs=14, gap=10)
footer(s, 4)

# =============================================================================
# 5 — MOMENT 2 : what/why + focus lifecycle
# =============================================================================
s = slide(); header(s, "Moment 2 · Accessibility", "Delete → Escape → focus returns")
whatwhy(s,
        "Click Delete — the dialog opens with focus on Cancel. Press Escape — it closes AND focus returns to the Delete button.",
        "Keyboard + screen-reader correctness you can't fake in a mockup. The full WAI-ARIA modal pattern, hand-rolled.")
y = Inches(3.75)
steps = [("Open", ["capture", "trigger el"]), ("Focus Cancel", ["safe default"]),
         ("Tab trap", ["stays inside"]), ("Escape", ["closes"]), ("Return focus", ["to trigger"])]
xs = [Inches(0.6) + Inches(2.46) * i for i in range(5)]
for (t, subs), x in zip(steps, xs):
    box(s, x, y, Inches(2.2), Inches(1.15), t, subs, edge=AMBER, fs_title=12.5)
for i in range(4):
    arrow(s, xs[i] + Inches(2.2), y + Inches(0.57), xs[i + 1], y + Inches(0.57), color=AMBER)
para(textbox(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(0.6)),
     "role=\"dialog\" · aria-modal=\"true\" · aria-labelledby → title. Cancel is the deliberate first focus so an accidental Enter can't confirm a delete.",
     13, MUTED, italic=True, first=True)
footer(s, 5)

# =============================================================================
# 6 — MOMENT 2 : code
# =============================================================================
s = slide(); header(s, "Moment 2 · Code", "The focus lifecycle (confirm-dialog.component.ts)")
code_panel(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(4.2),
           "// on open (ngOnChanges) — snapshot the trigger BEFORE the dialog steals focus\n"
           "this.previouslyFocused = document.activeElement as HTMLElement | null;\n"
           "this.needsInitialFocus = true;\n"
           "\n"
           "// after render (ngAfterViewChecked) — move focus to Cancel, deliberately\n"
           "if (this.needsInitialFocus && this.cancelBtn) this.cancelBtn.nativeElement.focus();\n"
           "\n"
           "// Escape closes\n"
           "@HostListener('document:keydown.escape') onEscape() { if (this.open) this.onCancel(); }\n"
           "\n"
           "// on close — return focus to exactly where the user was\n"
           "if (target?.focus) setTimeout(() => target.focus(), 0);", fs=12.5)
footer(s, 6)

# =============================================================================
# 7 — MOMENT 3 : Cypress
# =============================================================================
s = slide(); header(s, "Moment 3 · The test net", "Let the specs walk the UI")
panel(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.7), fill=CODE_BG, line=GREEN_DARK)
tf = textbox(s, Inches(0.85), Inches(1.83), Inches(11.6), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run(); r.text = "npm --prefix client run e2e:open"; _set(r, 15, CODE, MONO_FONT, bold=True)
# 4-spec table
hdr_y = Inches(2.65)
rows = [("employee-flow", "List, search, create + delete, required-field validation"),
        ("account-flow", "Account CRUD on detail page; balance validation; close + reopen"),
        ("audit-log-flow", "Each action lands the right entry at the top of the audit log"),
        ("employee-filters", "search/role/status/hasAccounts compose; reset; back to page 1")]
y = hdr_y
for i, (spec, desc) in enumerate(rows):
    panel(s, Inches(0.6), y, Inches(3.3), Inches(0.72), fill=PANEL2, line=ACCENT)
    tfa = textbox(s, Inches(0.75), y + Pt(6), Inches(3.0), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    ra = tfa.paragraphs[0].add_run(); ra.text = spec + ".cy.ts"; _set(ra, 12.5, GREEN_DEEP, MONO_FONT, bold=True)
    panel(s, Inches(4.0), y, Inches(8.7), Inches(0.72), fill=PANEL if i % 2 else WHITE, line=PANEL)
    tfb = textbox(s, Inches(4.2), y + Pt(6), Inches(8.4), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    rb = tfb.paragraphs[0].add_run(); rb.text = desc; _set(rb, 12, TEXT)
    y += Inches(0.82)
para(textbox(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.5)),
     "Pick one spec and stay quiet — watching a green run walk the real UI says more than narration.",
     13, AMBER, italic=True, first=True)
footer(s, 7)

# =============================================================================
# 8 — MOMENT 4 : what/why + flow
# =============================================================================
s = slide(); header(s, "Moment 4 · Tracing", "One id, browser → API → log")
whatwhy(s,
        "Open DevTools → Network → trigger any request. Show X-Correlation-Id on the request, then the same id in the server log.",
        "End-to-end traceability with zero extra dependencies — and it's OpenTelemetry-ready as-is.")
y = Inches(3.7)
nodes = [("Angular interceptor", ["mints", "X-Correlation-Id"]),
         ("Express middleware", ["echoes header", "+ req.correlationId"]),
         ("morgan logger", ["prints", "cid=…"])]
xs = [Inches(0.6), Inches(4.85), Inches(9.1)]
ws = [Inches(3.9), Inches(3.9), Inches(3.6)]
edges = [ACCENT, AMBER, AMBER]
for (t, subs), x, w, e in zip(nodes, xs, ws, edges):
    box(s, x, y, w, Inches(1.15), t, subs, edge=e, fs_title=13)
arrow(s, xs[0] + ws[0], y + Inches(0.57), xs[1], y + Inches(0.57))
arrow(s, xs[1] + ws[1], y + Inches(0.57), xs[2], y + Inches(0.57), color=AMBER)
para(textbox(s, Inches(0.6), Inches(5.25), Inches(12.1), Inches(0.6)),
     "The response echoes the id back too, so the client error interceptor logs the same cid on failures.",
     13, MUTED, italic=True, first=True)
footer(s, 8)

# =============================================================================
# 9 — MOMENT 4 : code
# =============================================================================
s = slide(); header(s, "Moment 4 · Code", "Three small pieces, one id")
code_panel(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(2.05),
           "// correlation-id.interceptor.ts (client)\n"
           "const cid = crypto.randomUUID();\n"
           "const cloned = req.clone({\n"
           "  setHeaders: { 'X-Correlation-Id': cid }\n"
           "});\n"
           "return next(cloned);", fs=11.5)
code_panel(s, Inches(6.85), Inches(1.7), Inches(5.85), Inches(2.05),
           "// correlation-id.js (server middleware)\n"
           "const incoming =\n"
           "  req.header(CORRELATION_HEADER) || uuid();\n"
           "res.setHeader(CORRELATION_HEADER, incoming);\n"
           "req.correlationId = incoming;\n"
           "next();", fs=11.5)
code_panel(s, Inches(0.6), Inches(3.95), Inches(12.1), Inches(1.5),
           "// logger.js — morgan token pulls the cid off req\n"
           "morgan.token('cid', (req) => req.correlationId);\n"
           "module.exports = morgan(':method :url :status :response-time ms - cid=:cid');", fs=12)
para(textbox(s, Inches(0.6), Inches(5.65), Inches(12.1), Inches(0.6)),
     "Client mints it · server echoes + attaches to req · logger stamps every line. Same id, end to end.",
     13, ACCENT, italic=True, first=True)
footer(s, 9)

# =============================================================================
# 10 — CLOSING
# =============================================================================
s = slide(bg=GREEN)
rib = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), SW, Pt(5))
rib.fill.solid(); rib.fill.fore_color.rgb = YELLOW; rib.line.fill.background(); rib.shadow.inherit = False
tf = textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.8))
para(tf, "Four moments, one theme", 36, WHITE, TITLE_FONT, bold=True, first=True)
tf2 = textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(3.6))
para(tf2, "Each shows something a screenshot can't:", 16, YELLOW, BODY_FONT, bold=True, first=True, space_after=12)
for q in ["Audit log — the compliance story, working live",
          "Dialog focus — accessibility you have to see to believe",
          "Cypress — the safety net, proving itself",
          "Correlation id — production-grade tracing, already wired"]:
    para(tf2, q, 15, WHITE, space_after=8)
para(tf2, "Show the behaviour first. Open the code only when they ask \"how?\"", 13, PANEL2, italic=True, space_after=0)

out = "docs/DEMO_WALKTHROUGH.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
