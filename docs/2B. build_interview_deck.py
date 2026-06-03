"""
Build INTERVIEW_PRESENTATION.pptx from the content in docs/INTERVIEW_PLAYBOOK.md.

Regenerate with:   py -3 docs/build_interview_deck.py
Requires:          py -3 -m pip install python-pptx

Theme: dark slate, green accent. 16:9. Shape-based architecture + flow diagrams,
real code snippets pulled from this project.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (TD Bank colours) -----------------------------------------------
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x00, 0x8A, 0x00)   # TD Green (primary)
GREEN_DARK = RGBColor(0x00, 0x6B, 0x00)   # hover/pressed
GREEN_DEEP = RGBColor(0x00, 0x4F, 0x00)   # text on light green
YELLOW     = RGBColor(0xFF, 0xC7, 0x2C)   # TD accent gold
BG      = WHITE
PANEL   = RGBColor(0xF2, 0xF4, 0xF5)      # light surface
PANEL2  = RGBColor(0xE8, 0xF5, 0xE8)      # TD green soft (tiles, alt rows)
ACCENT  = GREEN                            # primary accent
ACCENT2 = GREEN_DARK                       # secondary edges
AMBER   = RGBColor(0xB4, 0x53, 0x09)       # warn / backend edges
RED     = RGBColor(0xC7, 0x10, 0x2E)       # TD red (critical)
TEXT    = RGBColor(0x1C, 0x1C, 0x1C)       # high-contrast body
MUTED   = RGBColor(0x55, 0x5E, 0x66)       # captions
CODE    = RGBColor(0xE8, 0xF5, 0xE8)       # light text on dark code panel

EMOJI = None  # keep fonts simple/portable
TITLE_FONT = "Segoe UI Semibold"
BODY_FONT  = "Segoe UI"
MONO_FONT  = "Consolas"

SW, SH = Emu(12192000), Emu(6858000)  # exact 16:9 widescreen

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ---- helpers -----------------------------------------------------------------
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def _set(run, size, color, font=BODY_FONT, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    return tf


def para(tf, text, size, color, font=BODY_FONT, bold=False, italic=False,
         space_after=6, level=0, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.level = level
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    _set(r, size, color, font, bold, italic)
    return p


def accent_bar(s, x=Inches(0.6), y=Inches(0.62), w=Inches(0.9)):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(5))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False


def header(s, kicker, title):
    tf = textbox(s, Inches(0.6), Inches(0.30), Inches(12.1), Inches(0.3))
    para(tf, kicker.upper(), 12, ACCENT, BODY_FONT, bold=True, first=True, space_after=0)
    accent_bar(s)
    tf2 = textbox(s, Inches(0.6), Inches(0.74), Inches(12.1), Inches(0.9))
    para(tf2, title, 30, GREEN_DEEP, TITLE_FONT, bold=True, first=True, space_after=0)


def footer(s, n):
    tf = textbox(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Banking Admin Portal"
    _set(r, 9, MUTED)
    # page number right-aligned
    tfn = textbox(s, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.35))
    pn = tfn.paragraphs[0]; pn.alignment = PP_ALIGN.RIGHT
    rn = pn.add_run(); rn.text = str(n); _set(rn, 9, MUTED)


def panel(s, x, y, w, h, fill=PANEL, line=None, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(1.25)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def box(s, x, y, w, h, lines, fill=PANEL2, edge=ACCENT2, title_color=TEXT,
        body_color=MUTED, fs_title=13, fs_body=10):
    """lines: (title, [sub, sub]) ; draws a rounded box with title + sublines."""
    shp = panel(s, x, y, w, h, fill=fill, line=edge)
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    title, subs = lines
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; _set(r, fs_title, title_color, BODY_FONT, bold=True)
    for sub in subs:
        sp = tf.add_paragraph(); sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run(); sr.text = sub; _set(sr, fs_body, body_color, MONO_FONT)
    return shp


def arrow(s, x1, y1, x2, y2, color=ACCENT, width=2.0):
    c = s.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = straight
    c.line.color.rgb = color; c.line.width = Pt(width)
    c.shadow.inherit = False
    # arrowhead
    ln = c.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return c


def code_panel(s, x, y, w, h, code, fs=12):
    panel(s, x, y, w, h, fill=RGBColor(0x06, 0x2A, 0x14), line=GREEN_DARK)
    tf = textbox(s, x + Pt(10), y + Pt(8), w - Pt(20), h - Pt(16))
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0); p.line_spacing = 1.05
        r = p.add_run(); r.text = line if line else " "
        _set(r, fs, CODE, MONO_FONT)


def bullets(s, x, y, w, h, items, fs=15, gap=8):
    """items: list of (text, level, color, bold)."""
    tf = textbox(s, x, y, w, h)
    for i, it in enumerate(items):
        text, level, color, bold = (it + (0, TEXT, False))[:4] if isinstance(it, tuple) else (it, 0, TEXT, False)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level; p.space_after = Pt(gap)
        bullet = "▸ " if level == 0 else "– "
        r = p.add_run(); r.text = bullet + text
        _set(r, fs - (level * 1), color, BODY_FONT, bold=bold)
    return tf


def table(s, x, y, w, headers, rows, col_w=None, fs=12, header_fs=12, row_h=Inches(0.5)):
    nrows = len(rows) + 1; ncols = len(headers)
    h = row_h * nrows
    gtbl = s.shapes.add_table(nrows, ncols, x, y, w, h).table
    if col_w:
        for i, cw in enumerate(col_w):
            gtbl.columns[i].width = cw
    # style header
    for j, htext in enumerate(headers):
        c = gtbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT
        tf = c.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = htext
        _set(r, header_fs, BG, BODY_FONT, bold=True)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gtbl.cell(i + 1, j)
            c.fill.solid(); c.fill.fore_color.rgb = PANEL if i % 2 == 0 else PANEL2
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = val
            _set(r, fs, TEXT, BODY_FONT)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
    return gtbl


# =============================================================================
# 1 — TITLE
# =============================================================================
s = slide(bg=GREEN)
# gold ribbon
rib = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.55), SW, Pt(5))
rib.fill.solid(); rib.fill.fore_color.rgb = YELLOW; rib.line.fill.background(); rib.shadow.inherit = False
tf = textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.0))
para(tf, "BANKING ADMIN PORTAL", 16, YELLOW, BODY_FONT, bold=True, first=True, space_after=2)
tf2 = textbox(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.4))
para(tf2, "Managing employees & their linked accounts", 40, WHITE, TITLE_FONT, bold=True, first=True, space_after=4)
para(tf2, "An Angular 17 + NgRx admin tool, built as a take-home and treated as a real product", 18, PANEL2)
tf3 = textbox(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.0))
para(tf3, "Interview Walkthrough", 16, WHITE, BODY_FONT, bold=True, first=True, space_after=4)
para(tf3, "Angular 17 · TypeScript · NgRx + Signals · RxJS · Express · Jest · Cypress", 13, YELLOW, MONO_FONT)

# =============================================================================
# 2 — AGENDA
# =============================================================================
s = slide(); header(s, "Agenda", "What we'll walk through")
left = [
    ("The project in one minute", 0, TEXT, True),
    ("The 2-minute pitch", 0, TEXT, True),
    ("System architecture", 0, TEXT, True),
    ("NgRx + Signals — the split", 0, TEXT, True),
    ("Backend layering (MVC + repo + service)", 0, TEXT, True),
]
right = [
    ("Key decisions & trade-offs", 0, TEXT, True),
    ("Code that shows the craft", 0, TEXT, True),
    ("Quantified impact (honest numbers)", 0, TEXT, True),
    ("What didn't work / what's missing", 0, TEXT, True),
    ("Scaling 10x  ·  Q&A", 0, TEXT, True),
]
bullets(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(4.6), left, fs=18, gap=16)
bullets(s, Inches(6.9), Inches(1.9), Inches(6.0), Inches(4.6), right, fs=18, gap=16)
footer(s, 2)

# =============================================================================
# 3 — THE PROJECT
# =============================================================================
s = slide(); header(s, "Context", "The project in one minute")
bullets(s, Inches(0.7), Inches(1.8), Inches(7.3), Inches(4.6), [
    ("An internal Banking Admin Portal: a back-office admin manages employees and their linked banking accounts.", 0, TEXT, False),
    ("Full CRUD on both domains, all five HTTP verbs, NgRx, reactive forms, tests.", 0, TEXT, False),
    ("The interesting question: what would this look like if it shipped behind SSO at a bank tomorrow?", 0, ACCENT, True),
    ("Sole developer — frontend, backend mock, tests, docs, accessibility, responsive.", 0, MUTED, False),
], fs=16, gap=14)
# stat panel
panel(s, Inches(8.4), Inches(1.8), Inches(4.3), Inches(4.4), fill=PANEL)
tf = textbox(s, Inches(8.7), Inches(2.0), Inches(3.7), Inches(4.0))
para(tf, "AT A GLANCE", 12, ACCENT, BODY_FONT, bold=True, first=True, space_after=10)
for k, v in [("Required features", "100%"), ("Bonus items", "3 / 3"),
             ("Frontend unit specs", "4"), ("Backend Jest suites", "2"),
             ("Cypress e2e flows", "4"), ("Responsive tiers", "phone → 4K")]:
    p = tf.add_paragraph(); p.space_after = Pt(8)
    r = p.add_run(); r.text = v + "  "; _set(r, 16, ACCENT2, BODY_FONT, bold=True)
    r2 = p.add_run(); r2.text = k; _set(r2, 12, MUTED)
footer(s, 3)

# =============================================================================
# 4 — THE PITCH
# =============================================================================
s = slide(); header(s, "The 2-minute pitch", "Context · Problem · Role")
items = [
    ("CONTEXT", "A take-home for a Senior Angular role at a Canadian bank — an internal portal to manage employees and accounts."),
    ("PROBLEM", "The brief was generic. The real question: build it like it's going to production at a bank, not like a spec speed-run."),
    ("ROLE", "Sole developer — scaffolding through ongoing refactors, across frontend, backend mock, tests, docs and accessibility."),
]
y = Inches(1.9)
for tag, body in items:
    panel(s, Inches(0.7), y, Inches(12.0), Inches(1.25), fill=PANEL)
    tf = textbox(s, Inches(1.0), y + Pt(8), Inches(11.4), Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, tag, 13, ACCENT, BODY_FONT, bold=True, first=True, space_after=3)
    para(tf, body, 15, TEXT)
    y += Inches(1.45)
footer(s, 4)

# =============================================================================
# 5 — THREE KEY DECISIONS
# =============================================================================
s = slide(); header(s, "The pitch, continued", "Three key decisions")
cards = [
    ("1 · Layered backend", ["Strict MVC + repository +", "service — not a single", "server.js. Business rules", "stay testable in isolation."]),
    ("2 · NgRx + Signals", ["NgRx stays the source of", "truth; Signals appear at", "three seams (toSignal,", "signal, effect)."]),
    ("3 · Append-only audit", ["Per-employee audit log", "with field-level diffs.", "Banks audit everything —", "the obvious extension."]),
]
x = Inches(0.7)
for title, lines in cards:
    panel(s, x, Inches(1.9), Inches(3.95), Inches(3.7), fill=PANEL, line=ACCENT2)
    tf = textbox(s, x + Pt(14), Inches(2.1), Inches(3.6), Inches(3.4))
    para(tf, title, 18, ACCENT, TITLE_FONT, bold=True, first=True, space_after=10)
    for ln in lines:
        para(tf, ln, 14, TEXT, space_after=4)
    x += Inches(4.1)
tf = textbox(s, Inches(0.7), Inches(5.9), Inches(12.0), Inches(0.8))
para(tf, "Impact: every required feature + all 3 bonus items · backend Jest on the two highest-risk services · 4 Cypress flows · responsive to 4K · WCAG 2.1 AA on interactive surfaces.", 13, MUTED, italic=True, first=True)
footer(s, 5)

# =============================================================================
# 6 — SYSTEM ARCHITECTURE (shape diagram)
# =============================================================================
s = slide(); header(s, "Architecture", "System overview")
# Browser
box(s, Inches(5.3), Inches(1.55), Inches(2.7), Inches(0.55), ("Browser", []), fill=PANEL, edge=MUTED, fs_title=13)
# Angular panel
panel(s, Inches(0.7), Inches(2.35), Inches(12.0), Inches(1.95), fill=PANEL)
tf = textbox(s, Inches(0.85), Inches(2.4), Inches(3.0), Inches(0.4))
para(tf, "ANGULAR 17  ·  :4200", 12, ACCENT, BODY_FONT, bold=True, first=True)
box(s, Inches(0.95), Inches(2.85), Inches(2.5), Inches(1.2), ("Components", ["templates", "OnPush"]), edge=ACCENT2)
box(s, Inches(3.75), Inches(2.85), Inches(2.5), Inches(1.2), ("Facade", ["Employee/", "Account"]), edge=ACCENT2)
box(s, Inches(6.55), Inches(2.85), Inches(2.5), Inches(1.2), ("NgRx store", ["actions/reducer", "selectors/effects"]), edge=ACCENT2)
box(s, Inches(9.35), Inches(2.85), Inches(3.0), Inches(1.2), ("Interceptors", ["correlation-id", "error → ApiError"]), edge=ACCENT2)
arrow(s, Inches(3.45), Inches(3.45), Inches(3.75), Inches(3.45))
arrow(s, Inches(6.25), Inches(3.45), Inches(6.55), Inches(3.45))
arrow(s, Inches(9.05), Inches(3.45), Inches(9.35), Inches(3.45))
# proxy arrow
arrow(s, Inches(6.65), Inches(4.30), Inches(6.65), Inches(4.75), color=ACCENT)
tf = textbox(s, Inches(6.8), Inches(4.32), Inches(3.0), Inches(0.4))
para(tf, "proxy /api/*", 11, MUTED, MONO_FONT, first=True)
# Express panel
panel(s, Inches(0.7), Inches(4.8), Inches(12.0), Inches(1.95), fill=PANEL)
tf = textbox(s, Inches(0.85), Inches(4.85), Inches(4.0), Inches(0.4))
para(tf, "EXPRESS MOCK API  ·  :3000", 12, ACCENT, BODY_FONT, bold=True, first=True)
labels = [("Routes", []), ("Controllers", ["HTTP only"]), ("Services", ["rules + audit"]),
          ("Repositories", ["pure CRUD"]), ("Store", ["in-memory"])]
x = Inches(0.95)
xs = []
for title, subs in labels:
    box(s, x, Inches(5.3), Inches(2.18), Inches(1.2), (title, subs), edge=AMBER, fs_title=12)
    xs.append(x)
    x += Inches(2.35)
for i in range(len(xs) - 1):
    arrow(s, xs[i] + Inches(2.18), Inches(5.9), xs[i + 1], Inches(5.9), color=AMBER)
footer(s, 6)

# =============================================================================
# 7 — NgRx + Signals split
# =============================================================================
s = slide(); header(s, "State management", "NgRx + Signals — the split")
para(textbox(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(0.5)),
     "NgRx stays the single source of truth. Signals appear at three specific seams — not sprinkled everywhere.",
     15, TEXT, first=True)
table(s, Inches(0.7), Inches(2.4), Inches(12.0),
      ["Seam", "Where", "Why it's the right tool"],
      [["toSignal()", "list / detail components", "Read facade state synchronously in templates; store unchanged"],
       ["signal()", "UI flags (showForm, confirmOpen)", "Never shared/persisted — an NgRx slice would be ceremony"],
       ["computed()", "page summary, form heading", "Memoised derivations of multiple signals"],
       ["effect()", "document.title sync", "State-driven DOM side effect, auto-tracked + auto-disposed"]],
      col_w=[Inches(2.0), Inches(4.3), Inches(5.7)], fs=12, row_h=Inches(0.72))
para(textbox(s, Inches(0.7), Inches(6.3), Inches(12.0), Inches(0.6)),
     "\"I didn't replace NgRx with Signals — I bolted them on at the leaf.\"",
     15, ACCENT, italic=True, first=True)
footer(s, 7)

# =============================================================================
# 8 — NgRx unidirectional flow (shape diagram)
# =============================================================================
s = slide(); header(s, "Data flow", "NgRx — one direction")
y = Inches(2.4)
nodes = [("Component", Inches(0.7)), ("Effect", Inches(3.4)), ("API service", Inches(6.1)), ("Express", Inches(8.8))]
for name, x in nodes:
    box(s, x, y, Inches(2.2), Inches(0.95), (name, []), edge=ACCENT2, fs_title=14)
for i in range(3):
    arrow(s, nodes[i][1] + Inches(2.2), y + Inches(0.47), nodes[i + 1][1], y + Inches(0.47))
# labels
for txt, x in [("dispatch action", Inches(0.9)), ("calls", Inches(3.7)), ("HTTP", Inches(6.4))]:
    para(textbox(s, x, y - Inches(0.45), Inches(2.2), Inches(0.4)), txt, 10, MUTED, MONO_FONT, first=True)
# down to reducer/store and back
box(s, Inches(3.4), Inches(4.2), Inches(2.2), Inches(0.95), ("Reducer", []), edge=AMBER, fs_title=14)
box(s, Inches(6.1), Inches(4.2), Inches(2.2), Inches(0.95), ("Store", []), edge=AMBER, fs_title=14)
arrow(s, Inches(4.5), y + Inches(0.95), Inches(4.5), Inches(4.2), color=AMBER)
para(textbox(s, Inches(4.6), Inches(3.5), Inches(3.0), Inches(0.4)), "success/failure action", 10, MUTED, MONO_FONT, first=True)
arrow(s, Inches(5.6), Inches(4.67), Inches(6.1), Inches(4.67), color=AMBER)
# store -> facade -> back to component (selector)
box(s, Inches(0.7), Inches(4.2), Inches(2.2), Inches(0.95), ("Facade", ["selector"]), edge=ACCENT, fs_title=13)
arrow(s, Inches(6.1), Inches(4.67), Inches(2.9), Inches(4.67), color=ACCENT)
arrow(s, Inches(1.8), Inches(4.2), Inches(1.8), y + Inches(0.95), color=ACCENT)
para(textbox(s, Inches(0.7), Inches(5.35), Inches(8.0), Inches(0.5)),
     "Component dispatches; only the effect touches HTTP; only the reducer writes the store; the facade is the read/write interface.",
     13, MUTED, italic=True, first=True)
footer(s, 8)

# =============================================================================
# 9 — Backend layering
# =============================================================================
s = slide(); header(s, "Backend", "Layered MVC + repository + service")
code_panel(s, Inches(0.7), Inches(1.8), Inches(6.0), Inches(4.7),
"""server/
  routes/        URL -> controller
  controllers/   HTTP shape only
  services/      business rules
                 + audit + cascade
  repositories/  pure CRUD
  data/store.js  in-memory store
  validators/    errors[] arrays
  middleware/    cid, logger, errors
  utils/         problem-details,
                 sanitize""", fs=13)
bullets(s, Inches(7.0), Inches(1.9), Inches(5.7), Inches(4.5), [
    ("One-way dependency:", 0, ACCENT, True),
    ("controller → service → repository → store", 1, ACCENT2, False),
    ("Store is swappable — only repositories change when in-memory becomes Postgres.", 0, TEXT, False),
    ("Business rules unit-testable against a stub repository, no HTTP server needed.", 0, TEXT, False),
    ("Smell to watch: a controller importing a repository, or a service touching the store.", 0, MUTED, False),
], fs=14, gap=12)
footer(s, 9)

# =============================================================================
# 10 — Request lifecycle
# =============================================================================
s = slide(); header(s, "Backend", "Request lifecycle")
code_panel(s, Inches(0.7), Inches(1.8), Inches(12.0), Inches(4.7),
"""HTTP request
   |
   v
[middleware]   correlation-id  ->  logger
   |
   v
[routes]       /api/employees/:id   ->   controller.getById
   |
   v
[controller]   sanitize  ->  validate  ->  uniqueness check  ->  service call
   |
   v
[service]      apply business rules (defaults, timestamps, cascades, audit)
   |
   v
[repository]   reads / mutates  data/store.js
   |
   v
HTTP response  (or  [error-handler] -> RFC 7807 problem-details JSON)""", fs=13)
footer(s, 10)

# =============================================================================
# 11 — Decision tree / trade-offs
# =============================================================================
s = slide(); header(s, "Trade-offs", "Why this, over the alternatives")
table(s, Inches(0.7), Inches(1.9), Inches(12.0),
      ["Decision", "Chosen", "Rejected — and why"],
      [["State", "NgRx + Signals at the leaf", "NgRx-only (template ceremony); Signals-only (lose DevTools)"],
       ["Backend", "MVC + repo + service", "single server.js — no home for audit-diff logic"],
       ["Page size", "clamp at 100", "400 reject — breaks a caller who didn't know the limit"],
       ["Delete", "soft-close accts; hard+cascade employee", "hard-delete all — destroys audit history"],
       ["Validation", "hand-rolled functions", "zod/joi — extra dep; rules simple enough here"]],
      col_w=[Inches(1.7), Inches(4.0), Inches(6.3)], fs=11.5, row_h=Inches(0.72))
para(textbox(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.5)),
     "Every choice has the rejected alternative named — that's the senior signal.",
     13, ACCENT, italic=True, first=True)
footer(s, 11)

# =============================================================================
# 12 — CODE: cascade delete
# =============================================================================
s = slide(); header(s, "Code · backend hero", "Cascade delete + audit (employee.service.js)")
code_panel(s, Inches(0.7), Inches(1.75), Inches(12.0), Inches(4.0),
"""remove(employeeId, context) {
  const employee = EmployeeRepository.findById(employeeId);
  if (!employee) return false;

  // Snapshot OPEN accounts about to be cascade-closed.
  const toClose = AccountRepository.findByEmployeeId(employeeId)
    .filter((a) => a.status === 'OPEN').map((a) => ({ ...a }));

  EmployeeRepository.deleteById(employeeId);
  AccountRepository.updateAllByEmployeeId(employeeId, { status: 'CLOSED', updatedAt: nowIso() });

  AuditService.recordEmployeeDeleted(employee, context);
  for (const acc of toClose) AuditService.recordAccountCascadeClosed(acc, context);
  return true;
}""", fs=12.5)
para(textbox(s, Inches(0.7), Inches(5.95), Inches(12.0), Inches(0.8)),
     "One cohesive, transaction-shaped function orchestrating two repositories + the audit service. Has its own Jest suite — proves the cascade without booting Express.",
     13, MUTED, italic=True, first=True)
footer(s, 12)

# =============================================================================
# 13 — CODE: signals at the leaf
# =============================================================================
s = slide(); header(s, "Code · frontend seam", "Signals bridged onto NgRx (employee-detail)")
code_panel(s, Inches(0.7), Inches(1.75), Inches(12.0), Inches(3.6),
"""protected readonly employee =
  toSignal(this.facade.selected$, { initialValue: null });

constructor() {
  // state-driven DOM side effect: keep the tab title in sync
  effect(() => {
    const emp = this.employee();
    document.title = emp
      ? `${emp.firstName} ${emp.lastName} - ${DEFAULT_TITLE}`
      : DEFAULT_TITLE;
  });
}""", fs=13)
para(textbox(s, Inches(0.7), Inches(5.55), Inches(12.0), Inches(1.0)),
     "toSignal() bridges a facade observable to the template; effect() handles a DOM side effect, auto-tracking employee() and tearing itself down via DestroyRef. NgRx stays the source of truth.",
     13, MUTED, italic=True, first=True)
footer(s, 13)

# =============================================================================
# 14 — CODE: async validator
# =============================================================================
s = slide(); header(s, "Code · forms craft", "Async unique-email validator")
code_panel(s, Inches(0.7), Inches(1.75), Inches(12.0), Inches(3.7),
"""return of(value).pipe(
  debounceTime(150),        // don't hit the API on every keystroke
  distinctUntilChanged(),
  switchMap((email) => api.isEmailAvailable(email, excludeIdProvider()).pipe(
    map((available) => (available ? null : { emailTaken: true })),
    catchError(() => of(null))  // network error -> don't block; server re-checks
  )),
  first()
);""", fs=12.5)
para(textbox(s, Inches(0.7), Inches(5.65), Inches(12.0), Inches(1.0)),
     "A validator factory taking the API service + an excludeIdProvider (so an edit form's own email isn't a clash). Debounced, deduped, fails-open on network error, fires on blur.",
     13, MUTED, italic=True, first=True)
footer(s, 14)

# =============================================================================
# 15 — CODE: interceptors
# =============================================================================
s = slide(); header(s, "Code · HTTP", "Two functional interceptors, ordered on purpose")
code_panel(s, Inches(0.7), Inches(1.8), Inches(12.0), Inches(1.5),
"""provideHttpClient(
  withFetch(),
  withInterceptors([correlationIdInterceptor, errorInterceptor])
)""", fs=13)
bullets(s, Inches(0.7), Inches(3.6), Inches(12.0), Inches(2.6), [
    ("correlationIdInterceptor — attaches X-Correlation-Id per request (crypto.randomUUID()).", 0, TEXT, False),
    ("errorInterceptor — normalises HttpErrorResponse into ApiError, preserves the cid, logs it.", 0, TEXT, False),
    ("Order matters: cid runs first so the error handler can read the id it attached.", 0, ACCENT, True),
    ("Functional interceptors — the modern, tree-shakeable style, not class-based providers.", 0, MUTED, False),
], fs=14, gap=12)
footer(s, 15)

# =============================================================================
# 16 — Audit log
# =============================================================================
s = slide(); header(s, "Domain feature", "Append-only audit log")
bullets(s, Inches(0.7), Inches(1.85), Inches(6.0), Inches(4.5), [
    ("Every meaningful write is recorded as an immutable entry.", 0, TEXT, False),
    ("Per-action shape:", 0, ACCENT, True),
    ("CREATE / DELETE — full snapshot", 1, TEXT, False),
    ("UPDATE — field-level { before, after } diff", 1, TEXT, False),
    ("CLOSE / REOPEN / CASCADE_CLOSE — narrative", 1, TEXT, False),
    ("Attributed to actor + the request's correlation id.", 0, TEXT, False),
    ("Repository exposes no update/delete — append-only by design.", 0, MUTED, False),
], fs=14, gap=10)
code_panel(s, Inches(7.0), Inches(1.85), Inches(5.7), Inches(4.3),
"""Admin clicks
 "Mark INACTIVE"
   |
   v  PATCH /employees/:id
 service.patch()
   findById()  (BEFORE)
   update()    (write)
   audit.diff()
     { field: 'status',
       before: 'ACTIVE',
       after: 'INACTIVE' }
   |
   v  entry at top of log""", fs=12)
footer(s, 16)

# =============================================================================
# 17 — Quantified impact
# =============================================================================
s = slide(); header(s, "Outcomes", "Quantified impact — the honest numbers")
table(s, Inches(0.7), Inches(1.9), Inches(12.0),
      ["Metric", "Value"],
      [["Required spec features", "100%"],
       ["Listed bonus items", "3 of 3 (sort+paginate, audit log, Cypress e2e)"],
       ["Frontend unit specs", "4 — reducer, effect, service, form"],
       ["Backend Jest suites", "2 — the two highest-risk services (~30 cases)"],
       ["Cypress e2e flows", "4 — employee, account, audit, filters"],
       ["Responsive tiers", "phone ≤480 → tablet → desktop → 4K"],
       ["Accessibility", "WCAG 2.1 AA on interactive surfaces"]],
      col_w=[Inches(4.5), Inches(7.5)], fs=12.5, row_h=Inches(0.56))
para(textbox(s, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.5)),
     "Tests concentrated where a regression silently corrupts data — not spread thin for a vanity percentage.",
     13, MUTED, italic=True, first=True)
footer(s, 17)

# =============================================================================
# 18 — What didn't work
# =============================================================================
s = slide(); header(s, "Reflection", "The messy parts — what didn't work")
bullets(s, Inches(0.7), Inches(1.9), Inches(12.0), Inches(4.6), [
    ("Branded EmployeeId type — reverted: ~15-file cascade for a marginal safety win. Shipped a lighter type alias.", 0, TEXT, False),
    ("@Output('confirm') alias broke strict templates — renamed to confirmed/cancelled (style-guide form anyway).", 0, TEXT, False),
    ("Pipe inside (click) action = silent parser error pointing at the wrong line — fixed by migrating to a signal read.", 0, TEXT, False),
    ("Cypress audit specs carried state across runs — rewrote to assert on the newest entry + unique values.", 0, TEXT, False),
    ("\"production-shaped\", not \"production-ready\" — pre-built the disclaimer into the pitch.", 0, ACCENT, True),
], fs=15, gap=16)
footer(s, 18)

# =============================================================================
# 19 — Scaling 10x
# =============================================================================
s = slide(); header(s, "Scaling", "How would you scale this 10x?")
table(s, Inches(0.7), Inches(1.9), Inches(12.0),
      ["Bottleneck", "10x answer"],
      [["In-memory store", "Postgres + indexes on email, (employeeId,accountId), (employeeId, ts DESC)"],
       ["Audit log growth", "Already append-only → event sourcing; Redis read model for per-employee query"],
       ["Single process", "Shared DB → stateless → horizontal scale behind a load balancer"],
       ["List payload", "Server-side pagination + MAX_PAGE_SIZE clamp; CDN-cache with Vary"],
       ["No tracing", "X-Correlation-Id already there → OpenTelemetry, near-zero app change"],
       ["First-load bundle", "Already lazy-loaded per route; audit component ships with detail only"]],
      col_w=[Inches(2.9), Inches(9.1)], fs=11.5, row_h=Inches(0.62))
para(textbox(s, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.5)),
     "A shared database unlocks stateless horizontal scaling — naming that ordering is the senior move.",
     13, ACCENT, italic=True, first=True)
footer(s, 19)

# =============================================================================
# 20 — What's missing
# =============================================================================
s = slide(); header(s, "Self-review", "What's missing (I wrote my own PR review)")
# critical
para(textbox(s, Inches(0.7), Inches(1.8), Inches(12.0), Inches(0.35)), "CRITICAL — would block production", 13, RED, bold=True, first=True)
bullets(s, Inches(0.7), Inches(2.15), Inches(12.0), Inches(1.5), [
    ("No authentication — every audit actor is 'admin'. JWT middleware unblocks everything.", 0, TEXT, False),
    ("No direct tests on validators / controllers / middleware (covered only via Cypress).", 0, TEXT, False),
    ("Validation rules duplicated client + server — risk of silent contract drift.", 0, TEXT, False),
], fs=13, gap=7)
para(textbox(s, Inches(0.7), Inches(4.0), Inches(12.0), Inches(0.35)), "SHOULD FIX", 13, AMBER, bold=True, first=True)
bullets(s, Inches(0.7), Inches(4.35), Inches(12.0), Inches(1.4), [
    ("No CI pipeline · no rate limiting · prefers-reduced-motion not respected.", 0, TEXT, False),
    ("Confirm dialog rolls its own focus management instead of @angular/cdk/dialog.", 0, TEXT, False),
], fs=13, gap=7)
para(textbox(s, Inches(0.7), Inches(5.75), Inches(12.0), Inches(0.7)),
     "No code-review safety net on a solo build — so I wrote my own senior-dev review and shipped the top gaps in follow-up commits.",
     13, ACCENT, italic=True, first=True)
footer(s, 20)

# =============================================================================
# 21 — Closing / Q&A
# =============================================================================
s = slide(bg=GREEN)
rib = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), SW, Pt(5))
rib.fill.solid(); rib.fill.fore_color.rgb = YELLOW; rib.line.fill.background(); rib.shadow.inherit = False
tf = textbox(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.8))
para(tf, "Thank you — questions?", 38, WHITE, TITLE_FONT, bold=True, first=True)
tf2 = textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(3.5))
para(tf2, "Where would you like to go deeper?", 16, YELLOW, BODY_FONT, bold=True, first=True, space_after=14)
for q in ["The NgRx + Signals boundary",
          "The backend layering & cascade-delete transaction",
          "The audit-log design and the compliance angle",
          "The trade-offs I'd reverse (validation drift, auth)"]:
    para(tf2, q, 15, WHITE, space_after=8)
para(tf2, "Curiosity flip: \"What does a great senior-Angular demo look like to you — I'd rather hear what you wish more candidates did.\"",
     13, PANEL2, italic=True, space_after=0)

out = "docs/INTERVIEW_PRESENTATION.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
